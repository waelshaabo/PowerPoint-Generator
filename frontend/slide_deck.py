import os
import time
from pptx import Presentation
from pptx.util import Inches
import logging

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

def generate_unique_filename(base_name):
    timestamp = int(time.time())
    return f"{base_name}_{timestamp}.pptx"

class SlideDeck:

    def __init__(self, template_option, output_folder="generated"):
        self.template_option = template_option
        self.output_folder = output_folder
        self.prs = self._load_template()

    def _load_template(self):
        template_path = os.path.join("templates", f"{self.template_option}.pptx")
        logger.info(f"Attempting to load template: {template_path}")
        
        if os.path.exists(template_path):
            try:
                logger.info(f"Template file found: {template_path}")
                prs = Presentation(template_path)
                logger.info(f"Template loaded successfully: {self.template_option}")
                return prs
            except Exception as e:
                logger.error(f"Error loading template {self.template_option}: {str(e)}")
                logger.info("Falling back to blank presentation")
                return Presentation()
        else:
            logger.warning(f"Template file not found: {template_path}")
            logger.info("Using blank presentation")
            return Presentation()

    def add_slide(self, slide_data):
        prs = self.prs
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        # Title
        title_shape = shapes.title
        title_shape.text = slide_data.get("title_text", "")

        # Body
        if "text" in slide_data:
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in slide_data.get("text", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

                if "p1" in slide_data:
                    p = tf.add_paragraph()
                    p.text = slide_data.get("p1")
                    p.level = 1

        if "img_path" in slide_data:
            cur_left = 6
            for img_path in slide_data.get("img_path", []):
                top = Inches(2)
                left = Inches(cur_left)
                height = Inches(4)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                cur_left += 1

    def add_title_slide(self, title_page_data):
        # title slide
        prs = self.prs
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        if "title_text" in title_page_data:
            title.text = title_page_data.get("title_text")
        if "subtitle_text" in title_page_data:
            subtitle.text = title_page_data.get("subtitle_text")

    def create_presentation(self, title_slide_info, slide_pages_data=[]):
        try:
            base_name = title_slide_info.get("title_text", "presentation").\
                lower().replace(",", "").replace(" ", "-")
            file_name = generate_unique_filename(base_name)
            file_path = os.path.join(self.output_folder, file_name)
            
            self.add_title_slide(title_slide_info)
            for slide_data in slide_pages_data:
                self.add_slide(slide_data)

            # Try to save the file, if it fails, try an alternative location
            try:
                logger.info(f"Attempting to save presentation to: {file_path}")
                self.prs.save(file_path)
                logger.info(f"Presentation saved successfully to: {file_path}")
            except PermissionError:
                alternative_path = os.path.join(os.path.expanduser("~"), "Documents", file_name)
                logger.warning(f"Failed to save to {file_path}, trying alternative path: {alternative_path}")
                self.prs.save(alternative_path)
                logger.info(f"Presentation saved successfully to alternative path: {alternative_path}")
                file_path = alternative_path

            return file_path
        except Exception as e:
            logger.error(f"Error in create_presentation: {str(e)}")
            raise e